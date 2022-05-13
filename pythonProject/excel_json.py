import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import json

# from pptx.enum.text import PP_ALIGN
project_list = {}

wb = openpyxl.load_workbook('/home/fskhex/Desktop/project_line.xlsx')

# 获取表单

sh = wb['Sheet1']

rows = sh.max_row

columns = sh.max_column

for row in range(2, rows):
    project_line = sh.cell(row=row, column=1).value
    project_name = sh.cell(row=row, column=2).value
    owner = sh.cell(row=row, column=3).value
    repo = sh.cell(row=row, column=4).value
    owner_repo = {"owner": owner, "repo": repo}
    if project_line in project_list:
        project_name_dict = project_list[project_line]
        if project_name not in project_name_dict:
            project_name_dict[project_name] = []
        project_name_dict[project_name].append(owner_repo)
    else:
        project_list[project_line] = {}
        project_list[project_line][project_name] = []
        project_list[project_line][project_name].append(owner_repo)
print(project_list)

project_list = {}

wb = openpyxl.load_workbook('/home/fskhex/Desktop/project_line.xlsx')

# 获取表单

sh = wb['Sheet1']

rows = sh.max_row

columns = sh.max_column

for row in range(2, rows + 1):
    project_line = sh.cell(row=row, column=1).value
    project_name = sh.cell(row=row, column=2).value
    owner = sh.cell(row=row, column=3).value
    repo = sh.cell(row=row, column=4).value
    owner_repo = {"owner": owner, "repo": repo}
    if project_line in project_list:
        project_name_dict = project_list[project_line]
        if project_name not in project_name_dict:
            project_name_dict[project_name] = []
        project_name_dict[project_name].append(owner_repo)
    else:
        project_list[project_line] = {}
        project_list[project_line][project_name] = []
        project_list[project_line][project_name].append(owner_repo)
print(project_list)

with open('project_list.json', 'w', encoding='utf8') as file:
    json.dump(project_list, file, indent=4, ensure_ascii=False)
for project_line in project_list.keys():
    print(project_line)
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[6]

    unexists_image_files = []
    for project_name in project_list[project_line].keys():

        owner_repos = project_list[project_line][project_name]
        for owner_repo in owner_repos:
            owner = owner_repo["owner"]
            repo = owner_repo["repo"]
            print(f"开源项目名称：{project_name}({owner}/{repo})")
            image_file = f'/home/fskhex/Desktop/images/projectchart_imgs_png/{owner}___{repo}.png'
            if os.path.exists(image_file):
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes

                # title_shape = shapes.title
                # title_shape.text = f"开源项目名称：{project_name}({owner}/{repo})"
                textbox = shapes.add_textbox(left=Inches(1),
                                             top=Inches(1.5),
                                             width=Inches(9),
                                             height=Inches(2)
                                             )
                ## 向文本框加入文字
                tf = textbox.text_frame
                para = tf.add_paragraph()  # 添加段落
                para.text = f"开源项目名称：{project_name}({owner}/{repo})"
                # para.alignment = PP_ALIGN.CENTER  # 居中
                ## 设置字体
                font = para.font
                font.size = Pt(24)  # 大小
                font.name = '华文彩云'  # 字体
                font.bold = True  # 加粗
                font.italic = False  # 倾斜
                # font.color.rgb = RGBColor(225, 225, 0)  # 黄色

                # body_shape = shapes.placeholders[1]
                # body_shape.placeholder_format.type
                # body_shape.insert_picture('/home/fskhex/Desktop/images/pytorch___fairseq.jpg')

                shapes.add_picture(image_file=image_file,
                                   left=Inches(0.5),
                                   top=Inches(3),
                                   width=Inches(8.5)
                                   # ,height=Inches(4)
                                   )
            else:
                unexists_image_files.append(image_file)
    prs.save(f'{project_line}.pptx')
print("unexists_image_files::", unexists_image_files)
